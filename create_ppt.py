#!/usr/bin/env python3
"""
Aventra Consulting - Company Profile PowerPoint Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Define brand colors
PRIMARY_COLOR = RGBColor(0, 51, 102)      # Dark Blue
SECONDARY_COLOR = RGBColor(0, 102, 153)   # Medium Blue
ACCENT_COLOR = RGBColor(0, 153, 204)      # Light Blue
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_title_slide(prs):
    """Create title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Company name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AVENTRA CONSULTING"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Reimagining Business Through Digital Transformation"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Company Profile 2025"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_about_slide(prs):
    """Create About Us slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "About Aventra Consulting"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Main content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(2))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Founded in 2025, Aventra Consulting is among the largest IT and business consulting services firms in the world. We are insights-driven and outcomes-based to help accelerate returns on your investments."
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_GRAY
    
    p2 = tf.add_paragraph()
    p2.text = ""
    p2 = tf.add_paragraph()
    p2.text = "Our success is built on a unique member-owner culture where our professionals are also shareholders. This structure aligns our interests with those of our clients and ensures long-term commitment to quality and excellence."
    p2.font.size = Pt(18)
    p2.font.color.rgb = DARK_GRAY
    
    # Stats boxes
    stats = [
        ("100+", "Professionals\nWorldwide"),
        ("20+", "Locations\nGlobally"),
        ("15+", "Years of\nExperience"),
        ("40%", "Annual\nGrowth")
    ]
    
    start_x = 0.7
    for i, (number, label) in enumerate(stats):
        # Stat box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(start_x + i * 3.1), Inches(4.5),
            Inches(2.8), Inches(2.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = SECONDARY_COLOR
        box.line.fill.background()
        
        # Number
        num_box = slide.shapes.add_textbox(Inches(start_x + i * 3.1), Inches(4.7), Inches(2.8), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(start_x + i * 3.1), Inches(5.6), Inches(2.8), Inches(1))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_values_slide(prs):
    """Create Our Values slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Core Values"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    values = [
        ("Partnership", "We work side-by-side with our clients as trusted advisors, not just service providers."),
        ("Quality", "We deliver excellence in everything we do, ensuring sustainable outcomes for our clients."),
        ("Integrity", "We operate with transparency, honesty, and respect in all our relationships."),
        ("Innovation", "We continuously evolve our capabilities to help clients navigate changing technology landscapes.")
    ]
    
    colors = [PRIMARY_COLOR, SECONDARY_COLOR, ACCENT_COLOR, RgbColor(0, 76, 153)]
    
    for i, (value, desc) in enumerate(values):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.3
        y = 1.7 + row * 2.8
        
        # Value box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()
        
        # Value title
        val_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3), Inches(5.4), Inches(0.6))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Value description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.9), Inches(5.4), Inches(1.4))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
    
    return slide

def add_services_slide(prs):
    """Create Services slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Services"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "End-to-end services to meet the ever-evolving digital expectations of your customers"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    
    services = [
        "Business Consulting",
        "Enterprise Solutions",
        "Managed IT Services",
        "Artificial Intelligence",
        "Data Analytics",
        "Cloud & Hybrid IT",
        "Core Tech Expertise",
        "Staff Augmentation",
        "Business Process Services"
    ]
    
    for i, service in enumerate(services):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 2.2 + row * 1.7
        
        # Service box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4), Inches(1.4))
        box.fill.solid()
        if row == 0:
            box.fill.fore_color.rgb = PRIMARY_COLOR
        elif row == 1:
            box.fill.fore_color.rgb = SECONDARY_COLOR
        else:
            box.fill.fore_color.rgb = ACCENT_COLOR
        box.line.fill.background()
        
        # Service text
        svc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.4), Inches(3.6), Inches(0.8))
        tf = svc_box.text_frame
        p = tf.paragraphs[0]
        p.text = service
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_services_detail_slide(prs):
    """Create Services Detail slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Service Offerings - Details"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    services_detail = [
        ("Core Tech Expertise", "VDI, Virtual Machine, Cloud transformation, EUC services, Service desk, Workplace engineering"),
        ("Enterprise Solutions", "Digital Transformation, SaaS/PaaS, Enterprise Mobility, CRM Development, Salesforce"),
        ("IT Consulting", "Real Estate Services, GCC expansion, Technology Solutions, Legal Assistance, Co-Location"),
        ("Artificial Intelligence", "AI Development, Agentic AI, Machine Learning, Data Analytics, Process Automation"),
        ("Business Process Services", "Process automation, Workflow optimization, Cost reduction, Quality improvement"),
        ("Staff Augmentation", "FTE supplies, Work to contract, TNM Model, Offshore Models, Employee Charge back")
    ]
    
    for i, (title, details) in enumerate(services_detail):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.3
        y = 1.5 + row * 2
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(6), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = "▸ " + title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Details
        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.5), Inches(5.7), Inches(1.3))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = details
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
    
    return slide

def add_industries_slide(prs):
    """Create Industries slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Industries We Serve"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Industry-specific insights and solutions to navigate unique challenges and seize opportunities"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    
    industries = [
        ("Banking", "Modernize operations, enhance customer experiences"),
        ("Government", "Digital transformation for better citizen services"),
        ("Healthcare", "Innovative technology for patient outcomes"),
        ("Manufacturing", "Industry 4.0 and supply chain optimization"),
        ("Retail", "Omnichannel experiences for consumers"),
        ("Energy & Utilities", "Sustainable, digital-first operations"),
        ("Oil & Gas", "Safety, compliance, and efficiency"),
        ("IT", "Scalable digital infrastructure"),
        ("Insurance", "Risk management and compliance")
    ]
    
    for i, (industry, desc) in enumerate(industries):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 2.1 + row * 1.75
        
        # Industry box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = SECONDARY_COLOR
        
        # Industry name
        ind_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(3.6), Inches(0.5))
        tf = ind_box.text_frame
        p = tf.paragraphs[0]
        p.text = industry
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.65), Inches(3.6), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
    
    return slide

def add_geography_slide(prs):
    """Create Global Presence slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Global Presence"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "We provide a global antenna based on facts, not hype, to help improve returns on your IT and business investments"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    
    regions = [
        ("Middle East", "Saudi Arabia, Bahrain, UAE, Kuwait"),
        ("Europe", "UK, France, Germany, Sweden, Italy, Norway"),
        ("Australia & Oceania", "Australia, New Zealand"),
        ("Americas", "Columbia, Puerto Rico")
    ]
    
    colors = [PRIMARY_COLOR, SECONDARY_COLOR, ACCENT_COLOR, RgbColor(0, 76, 153)]
    
    for i, (region, countries) in enumerate(regions):
        x = 0.5 + i * 3.2
        
        # Region box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(3), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()
        
        # Region name
        reg_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.6), Inches(2.6), Inches(0.7))
        tf = reg_box.text_frame
        p = tf.paragraphs[0]
        p.text = region
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Countries
        ctry_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.5), Inches(2.6), Inches(3))
        tf = ctry_box.text_frame
        tf.word_wrap = True
        for country in countries.split(", "):
            p = tf.add_paragraph()
            p.text = "• " + country
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
    
    return slide

def add_leadership_slide(prs):
    """Create Leadership slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Leadership Team"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Decades of experience in technology, business consulting, and strategic management"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    
    leaders = [
        ("Vikas Grover", "President and CEO"),
        ("Deepinti Grover", "Chief Operating Officer"),
        ("Rovil Mahajan", "Executive Vice-President and CFO")
    ]
    
    for i, (name, title) in enumerate(leaders):
        x = 1.5 + i * 4
        
        # Profile circle placeholder
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.5), Inches(2.5), Inches(2.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = SECONDARY_COLOR
        circle.line.color.rgb = PRIMARY_COLOR
        
        # Initials
        init_box = slide.shapes.add_textbox(Inches(x), Inches(3.2), Inches(2.5), Inches(1))
        tf = init_box.text_frame
        p = tf.paragraphs[0]
        initials = "".join([n[0] for n in name.split()])
        p.text = initials
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Name
        name_box = slide.shapes.add_textbox(Inches(x - 0.5), Inches(5.2), Inches(3.5), Inches(0.6))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x - 0.5), Inches(5.7), Inches(3.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_case_study_slide(prs):
    """Create Case Study slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Excellence Delivered"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Case study box
    case_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.2))
    case_box.fill.solid()
    case_box.fill.fore_color.rgb = LIGHT_GRAY
    case_box.line.color.rgb = SECONDARY_COLOR
    
    # Case title
    cs_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.6))
    tf = cs_title.text_frame
    p = tf.paragraphs[0]
    p.text = "MAPFRE USA - Insurance Modernization"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Case description
    cs_desc = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1))
    tf = cs_desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "With our help, MAPFRE USA migrated to a modern core platform, driving faster claims, better underwriting and data-powered growth. Reimagining insurance through modernization."
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    
    # Value propositions
    values = [
        "Strategic business consulting to drive sustainable growth and organizational excellence",
        "End-to-end business process services enhancing efficiency, scalability, and operational performance",
        "Managed IT services to ensure optimal performance and security",
        "AI solutions to automate processes and enhance decision-making"
    ]
    
    for i, value in enumerate(values):
        y = 4.2 + i * 0.75
        val_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(0.7))
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "✓ " + value
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
    
    return slide

def add_contact_slide(prs):
    """Create Contact slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Let's Connect"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ready to transform your business?"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Contact card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3), Inches(8.3), Inches(3.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()
    
    # Address
    addr_box = slide.shapes.add_textbox(Inches(3), Inches(3.3), Inches(7.3), Inches(1.2))
    tf = addr_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📍 Global Headquarters"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    p2 = tf.add_paragraph()
    p2.text = "Akasa Coworking, 3rd Floor, Tower - B"
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_GRAY
    
    p3 = tf.add_paragraph()
    p3.text = "UNITECH CYBER PARK, Sector - 39, Gurugram, India"
    p3.font.size = Pt(14)
    p3.font.color.rgb = DARK_GRAY
    
    # Phone
    phone_box = slide.shapes.add_textbox(Inches(3), Inches(4.8), Inches(7.3), Inches(0.7))
    tf = phone_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📞 +91 9070030003 | +91 8713001002"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    
    # Email
    email_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7.3), Inches(0.7))
    tf = email_box.text_frame
    p = tf.paragraphs[0]
    p.text = "✉️ sales@consultaventra.com"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    
    # Website
    web_box = slide.shapes.add_textbox(Inches(3), Inches(6.2), Inches(7.3), Inches(0.5))
    tf = web_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🌐 www.consultaventra.com"
    p.font.size = Pt(16)
    p.font.color.rgb = SECONDARY_COLOR
    
    return slide

def add_thank_you_slide(prs):
    """Create Thank You slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Company name
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(1))
    tf = company_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AVENTRA CONSULTING"
    p.font.size = Pt(36)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(0.6))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Your Partner in Digital Transformation"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Main function to create the presentation"""
    # Create presentation with widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    add_title_slide(prs)
    add_about_slide(prs)
    add_values_slide(prs)
    add_services_slide(prs)
    add_services_detail_slide(prs)
    add_industries_slide(prs)
    add_geography_slide(prs)
    add_leadership_slide(prs)
    add_case_study_slide(prs)
    add_contact_slide(prs)
    add_thank_you_slide(prs)
    
    # Save the presentation
    output_path = "/app/Aventra_Consulting_Company_Profile.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved successfully to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
